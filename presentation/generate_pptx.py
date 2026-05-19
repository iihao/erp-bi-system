"""
Generate ERP-BI Defense PPTX
Usage: python3 generate_pptx.py
Output: erp-bi-system/presentation/ERP-BI 答辩_new.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Color Palette =====
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
PRIMARY_BLUE = RGBColor(0x25, 0x63, 0xEB)
LIGHT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xF9, 0xFA)
MEDIUM_GRAY = RGBColor(0x6B, 0x72, 0x80)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
VERY_LIGHT_BLUE = RGBColor(0xEE, 0xF2, 0xFF)
SLIDE_BG = RGBColor(0xF9, 0xFA, 0xFB)
TABLE_HEADER_BG = RGBColor(0x25, 0x63, 0xEB)
TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xFF)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS = os.path.join(BASE_DIR, 'screenshots')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg_gradient(slide, color1=DARK_BLUE, color2=RGBColor(0x0F, 0x17, 0x2A)):
    """Add solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_solid_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text, font_size=16, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(4), font_name='Microsoft YaHei'):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_icon_circle(slide, left, top, size, text, fill_color=PRIMARY_BLUE, text_color=WHITE, font_size=20):
    """Add a circle icon with text inside."""
    shape = add_shape(slide, left, top, size, size, fill_color=fill_color, shape_type=MSO_SHAPE.OVAL)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    # Center text vertically
    shape.text_frame.auto_size = None
    return shape


def add_card(slide, left, top, width, height, title, items, icon_text="", icon_color=PRIMARY_BLUE, title_color=DARK_BLUE, bg_color=WHITE, border_color=None):
    """Add a card-style container with title and bullet items."""
    # Card background
    card = add_shape(slide, left, top, width, height, fill_color=bg_color,
                     line_color=border_color if border_color else RGBColor(0xE5, 0xE7, 0xEB), line_width=Pt(1))

    # Icon circle
    if icon_text:
        add_icon_circle(slide, left + Inches(0.3), top + Inches(0.25), Inches(0.5), icon_text, fill_color=icon_color, font_size=18)

    # Title
    add_textbox(slide, left + Inches(0.9), top + Inches(0.25), width - Inches(1.2), Inches(0.45),
                title, font_size=16, bold=True, color=title_color)

    # Items
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.85), width - Inches(0.6), height - Inches(1.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_GRAY
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(4)
        p.space_after = Pt(4)


def add_table(slide, left, top, width, height, data, col_widths=None, header_color=TABLE_HEADER_BG):
    """Add a styled table. data is list of lists, first row = header."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = 'Microsoft YaHei'
                paragraph.alignment = PP_ALIGN.CENTER
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


# =============================================================================
# SLIDE 1: Cover
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg_gradient(slide)

# Decorative top accent line
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

# Title area
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.2),
            "基于 ERP 系统的商业智能报表的设计与实现",
            font_size=34, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.6),
            "毕业设计答辩",
            font_size=22, bold=False, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Decorative line
add_shape(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04), fill_color=ACCENT_CYAN)

# Info
info_items = [
    ("姓名：_______________", Inches(4.2)),
    ("导师：_______________", Inches(5.0)),
    ("日期：2026 年 3 月", Inches(5.8)),
    ("学校：_______________", Inches(6.6)),
]
for text, y in info_items:
    add_textbox(slide, Inches(1.5), y, Inches(10.3), Inches(0.5),
                text, font_size=16, color=RGBColor(0xCC, 0xD6, 0xDD), alignment=PP_ALIGN.CENTER)

# Bottom accent
add_shape(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), fill_color=PRIMARY_BLUE)


# =============================================================================
# SLIDE 2: Table of Contents
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

# Left sidebar
add_shape(slide, Inches(0), Inches(0), Inches(3.5), prs.slide_height, fill_color=DARK_BLUE)

# Title on sidebar
add_textbox(slide, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.8),
            "目录", font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_shape(slide, Inches(0.5), Inches(2.4), Inches(1.5), Inches(0.04), fill_color=ACCENT_CYAN)

# TOC items
toc_items = [
    ("01", "项目背景与目标"),
    ("02", "系统架构设计"),
    ("03", "数仓分层设计"),
    ("04", "ETL 流程实现"),
    ("05", "技术栈介绍"),
    ("06", "核心功能展示"),
    ("07", "AI 智能问数创新"),
    ("08", "系统测试与部署"),
    ("09", "总结与展望"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(3.0) + Inches(i * 0.48)
    add_textbox(slide, Inches(0.5), y, Inches(0.6), Inches(0.4),
                num, font_size=18, bold=True, color=ACCENT_CYAN)
    add_textbox(slide, Inches(1.2), y, Inches(2.0), Inches(0.4),
                title, font_size=16, color=RGBColor(0xE5, 0xE7, 0xEB))

# Right side content - key highlights
highlights = [
    ("🏗️", "数仓分层架构", "ODS → DWD → DWS → ADS 四层设计，共 24 张表"),
    ("🤖", "AI 智能问数", "自然语言转 SQL，准确率 > 80%"),
    ("💡", "双 BI 方案", "自研前端 + Metabase，优势互补"),
    ("🚀", "一键部署", "Docker Compose 编排，快速上线"),
]

for i, (icon, title, desc) in enumerate(highlights):
    y = Inches(1.2) + Inches(i * 1.3)
    # Card
    add_card(slide, Inches(4.5), y, Inches(8.0), Inches(1.1),
             title, [desc], icon_text=icon, icon_color=PRIMARY_BLUE,
             bg_color=WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB))


# =============================================================================
# SLIDE 3: Project Background
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

# Section header
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "01  项目背景与目标", font_size=28, bold=True, color=WHITE)

# Current situation card
add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5),
         "企业信息化建设现状",
         [
             "企业信息化建设深入，ERP 系统积累大量业务数据",
             "如何将数据转化为商业智能，支持管理决策",
             "传统 BI 系统门槛高，业务人员难以自助分析",
         ],
         icon_text="📊", icon_color=ACCENT_ORANGE)

# Goal card
add_card(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.7),
         "项目建设目标：构建低门槛、智能化的 BI 报表平台",
         [
             "让业务人员能够轻松获取数据洞察",
             "让管理层能够快速做出数据驱动决策",
             "让技术人员能够高效维护数据管道",
         ],
         icon_text="🎯", icon_color=PRIMARY_BLUE)

# Right side - visual emphasis
add_shape(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(7.7), Inches(2.0), Inches(4.5), Inches(0.8),
            "核心价值", font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.2), Inches(2.9), Inches(1.5), Inches(0.04), fill_color=ACCENT_CYAN)

value_items = [
    ("数据驱动", "从经验决策\n转向数据决策"),
    ("降本增效", "减少人工报表\n开发成本"),
    ("敏捷响应", "快速获取\n业务洞察"),
]
for i, (title, desc) in enumerate(value_items):
    y = Inches(3.3) + Inches(i * 1.4)
    add_icon_circle(slide, Inches(8.0), y, Inches(0.6), str(i + 1), fill_color=ACCENT_CYAN, font_size=16)
    add_textbox(slide, Inches(8.8), y - Inches(0.05), Inches(3.5), Inches(0.4),
                title, font_size=18, bold=True, color=WHITE)
    add_textbox(slide, Inches(8.8), y + Inches(0.35), Inches(3.5), Inches(0.6),
                desc, font_size=13, color=RGBColor(0x94, 0xA3, 0xB8))


# =============================================================================
# SLIDE 4: Project Goals & Innovation
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "02  项目目标与创新亮点", font_size=28, bold=True, color=WHITE)

# Core goals
add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.5),
         "核心目标",
         [
             "构建数仓分层架构（ODS → DWD → DWS → ADS）",
             "实现 ETL 数据流水线（自动化同步与转换）",
             "提供 BI 可视化能力（Metabase 集成）",
             "实现 AI 智能问数（自然语言查询）",
             "一键部署（Docker Compose）",
         ],
         icon_text="✅", icon_color=ACCENT_GREEN,
         title_color=DARK_BLUE, bg_color=WHITE, border_color=RGBColor(0xE5, 0xE7, 0xEB))

# Innovation highlights
add_card(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.5),
         "创新亮点",
         [
             "AI 智能问数 — 自然语言查询数据，准确率 > 80%",
             "双 BI 方案 — 自研前端 + Metabase 优势互补",
             "经典数仓四层架构 — 数据解耦，易于扩展",
             "Docker 容器化 — 5 个服务一键启动",
             "增量同步机制 — 基于时间戳，避免重复处理",
         ],
         icon_text="⭐", icon_color=ACCENT_ORANGE,
         title_color=DARK_BLUE, bg_color=VERY_LIGHT_BLUE, border_color=PRIMARY_BLUE)


# =============================================================================
# SLIDE 5: System Architecture
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "03  系统架构设计", font_size=28, bold=True, color=WHITE)

# Try to add architecture SVG as image (SVG not supported, use PNG if exists)
arch_path = os.path.join(BASE_DIR, 'architecture.png')
if os.path.exists(arch_path):
    slide.shapes.add_picture(arch_path, Inches(0.4), Inches(1.3), Inches(7.5), Inches(5.5))

# Architecture layers description on the right
layer_cards = [
    ("🖥️", "用户层", ["Vue3 前端：仪表板/报表/AI 问数", "Metabase BI：自助分析/拖拽报表", "REST API：外部系统集成"], LIGHT_BLUE),
    ("⚙️", "应用层", ["FastAPI 高性能后端", "AI 问数服务（百炼 Qwen3.5）", "API 网关 / 认证 / 查询引擎"], ACCENT_ORANGE),
    ("💾", "数据层", ["MySQL 8.0 业务数据库", "Redis 7.x 缓存加速", "数仓分层 ODS→DWD→DWS→ADS"], ACCENT_GREEN),
]

for i, (icon, title, items, color) in enumerate(layer_cards):
    y = Inches(1.3) + Inches(i * 1.9)
    add_card(slide, Inches(8.2), y, Inches(4.6), Inches(1.7),
             title, items, icon_text=icon, icon_color=color,
             bg_color=WHITE, border_color=color)


# =============================================================================
# SLIDE 6: Data Warehouse Layers
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "04  数仓分层设计", font_size=28, bold=True, color=WHITE)

# Table for warehouse layers
table_data = [
    ["层级", "名称", "描述", "表数量"],
    ["ODS", "原始数据层", "与业务库结构一致，直接抽取", "5 表"],
    ["DWD", "清洗标准化层", "统一格式、去重清洗、维度标准化", "5 表"],
    ["DWS", "轻度聚合层", "主题汇总、中度聚合", "6 表"],
    ["ADS", "报表指标层", "直接用于报表展示的预计算指标", "8 表"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(8.0), Inches(2.5), table_data,
          col_widths=[Inches(1.2), Inches(2.0), Inches(3.8), Inches(1.0)])

# Data flow
add_textbox(slide, Inches(0.6), Inches(4.3), Inches(8), Inches(0.5),
            "数据流转路径", font_size=18, bold=True, color=DARK_BLUE)

# Flow boxes
flow_items = [
    ("ERP 业务库", DARK_BLUE),
    ("ODS", LIGHT_BLUE),
    ("DWD", ACCENT_CYAN),
    ("DWS", ACCENT_GREEN),
    ("ADS", ACCENT_ORANGE),
    ("BI 展示", ACCENT_PURPLE),
]
for i, (text, color) in enumerate(flow_items):
    x = Inches(0.6) + Inches(i * 1.4)
    add_shape(slide, x, Inches(4.9), Inches(1.2), Inches(0.5), fill_color=color)
    add_textbox(slide, x, Inches(4.92), Inches(1.2), Inches(0.45),
                text, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_textbox(slide, x + Inches(1.2), Inches(4.92), Inches(0.2), Inches(0.45),
                    "→", font_size=14, bold=True, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Layer detail cards on right
layer_details = [
    ("ODS 层", ["ods_products - 产品信息", "ods_customers - 客户信息", "ods_orders - 订单主表", "ods_order_items - 订单明细", "ods_suppliers - 供应商信息"], LIGHT_BLUE),
    ("DWS 层", ["dws_sales_daily - 销售日报", "dws_sales_monthly - 销售月报", "dws_product_summary - 产品汇总", "dws_customer_stats - 客户统计"], ACCENT_CYAN),
]

for i, (title, items, color) in enumerate(layer_details):
    y = Inches(5.7) + Inches(i * 0.0)  # will adjust
    y = Inches(5.7) if i == 0 else Inches(6.65)
    add_card(slide, Inches(0.6), y, Inches(5.5), Inches(0.85),
             title, items[:2], icon_text="", icon_color=color,
             bg_color=WHITE, border_color=color)

# Right side: summary
add_shape(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(9.3), Inches(1.8), Inches(3.2), Inches(0.5),
            "数仓统计", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.8), Inches(2.4), Inches(2), Inches(0.04), fill_color=ACCENT_CYAN)

stats = [
    ("总表数", "24 张"),
    ("层级数", "4 层"),
    ("数据量", "~5,000 条"),
    ("同步耗时", "< 20 秒"),
    ("增量更新", "基于时间戳"),
]
for i, (label, value) in enumerate(stats):
    y = Inches(2.8) + Inches(i * 0.75)
    add_textbox(slide, Inches(9.3), y, Inches(3.2), Inches(0.35),
                label, font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9.3), y + Inches(0.35), Inches(3.2), Inches(0.35),
                value, font_size=22, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 7: ETL Process
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "05  ETL 流程实现", font_size=28, bold=True, color=WHITE)

# ETL flow diagram as boxes
etl_stages = [
    ("ODS 抽取", "全量/增量同步", LIGHT_BLUE),
    ("DWD 清洗", "去重/标准化", ACCENT_CYAN),
    ("DWS 聚合", "日报/月报", ACCENT_GREEN),
    ("ADS 报表", "KPI/排行", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(etl_stages):
    x = Inches(0.6) + Inches(i * 3.1)
    add_shape(slide, x, Inches(1.5), Inches(2.6), Inches(1.3), fill_color=color)
    add_textbox(slide, x, Inches(1.6), Inches(2.6), Inches(0.5),
                title, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.15), Inches(2.6), Inches(0.5),
                desc, font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, x + Inches(2.6), Inches(1.8), Inches(0.5), Inches(0.5),
                    "→", font_size=28, bold=True, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Core ETL scripts detail
add_textbox(slide, Inches(0.6), Inches(3.2), Inches(5), Inches(0.5),
            "核心 ETL 处理", font_size=20, bold=True, color=DARK_BLUE)

etl_details = [
    ("ODS 抽取", "全量同步 + 增量更新（基于 updated_at 时间戳）", PRIMARY_BLUE),
    ("DWD 清洗", "去重、空值处理、格式统一（日期、金额）", ACCENT_CYAN),
    ("DWS 聚合", "按日/月/产品/客户维度汇总", ACCENT_GREEN),
    ("ADS 报表", "预计算 KPI、排行、占比指标", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(etl_details):
    y = Inches(3.8) + Inches(i * 0.8)
    add_icon_circle(slide, Inches(0.7), y, Inches(0.45), str(i + 1), fill_color=color, font_size=14)
    add_textbox(slide, Inches(1.3), y - Inches(0.02), Inches(2), Inches(0.35),
                title, font_size=15, bold=True, color=color)
    add_textbox(slide, Inches(1.3), y + Inches(0.3), Inches(5), Inches(0.35),
                desc, font_size=13, color=DARK_GRAY)

# Data sync results table
add_textbox(slide, Inches(7.2), Inches(3.2), Inches(5), Inches(0.5),
            "数据同步效果", font_size=20, bold=True, color=DARK_BLUE)

sync_table = [
    ["层级", "数据量", "同步时间", "说明"],
    ["ODS", "~5,000 条", "< 10 秒", "从 ERP 业务库抽取"],
    ["DWD", "~4,800 条", "< 5 秒", "清洗后数据"],
    ["DWS", "~365 条", "< 3 秒", "按日聚合"],
    ["ADS", "~50 条", "< 2 秒", "预计算指标"],
]
add_table(slide, Inches(7.2), Inches(3.8), Inches(5.5), Inches(2.5), sync_table,
          col_widths=[Inches(1.0), Inches(1.2), Inches(1.0), Inches(2.3)])

# Incremental update
add_card(slide, Inches(7.2), Inches(6.5), Inches(5.5), Inches(0.7),
         "增量更新机制：基于 updated_at 时间戳，只同步变更数据",
         [], icon_text="🔄", icon_color=ACCENT_GREEN,
         bg_color=VERY_LIGHT_BLUE, border_color=ACCENT_GREEN)


# =============================================================================
# SLIDE 8: Tech Stack
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "06  技术栈介绍", font_size=28, bold=True, color=WHITE)

tech_data = [
    ["层级", "技术", "版本", "说明"],
    ["前端", "Vue3 + Vite + Element Plus", "3.4+", "响应式 UI 框架"],
    ["图表", "ECharts", "5.x", "数据可视化"],
    ["后端", "FastAPI", "0.100+", "高性能异步 API"],
    ["数据库", "MySQL", "8.0+", "关系型数据库"],
    ["缓存", "Redis", "7.x", "内存缓存加速"],
    ["BI", "Metabase", "0.47+", "开源 BI 工具"],
    ["AI", "百炼 Qwen3.5-Plus", "-", "自然语言处理"],
    ["容器", "Docker", "20.x", "容器化部署"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(8.0), Inches(4.5), tech_data,
          col_widths=[Inches(1.2), Inches(3.0), Inches(1.5), Inches(2.3)])

# Architecture diagram on right side
add_shape(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.5), fill_color=DARK_BLUE)
add_textbox(slide, Inches(9.3), Inches(1.8), Inches(3.2), Inches(0.5),
            "技术架构概览", font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.8), Inches(2.4), Inches(1.5), Inches(0.04), fill_color=ACCENT_CYAN)

tech_layers = [
    ("用户层", "Vue3 / ECharts / Metabase", LIGHT_BLUE),
    ("网关层", "FastAPI / Redis", ACCENT_ORANGE),
    ("数据层", "MySQL / 数仓四层", ACCENT_GREEN),
    ("AI 层", "百炼 Qwen3.5-Plus", ACCENT_PURPLE),
    ("部署层", "Docker Compose", ACCENT_CYAN),
]
for i, (layer, techs, color) in enumerate(tech_layers):
    y = Inches(2.8) + Inches(i * 0.85)
    add_icon_circle(slide, Inches(9.3), y, Inches(0.4), str(i + 1), fill_color=color, font_size=12)
    add_textbox(slide, Inches(9.9), y - Inches(0.02), Inches(2.5), Inches(0.3),
                layer, font_size=14, bold=True, color=WHITE)
    add_textbox(slide, Inches(9.9), y + Inches(0.3), Inches(2.5), Inches(0.3),
                techs, font_size=12, color=RGBColor(0x94, 0xA3, 0xB8))

# Docker compose at bottom
add_card(slide, Inches(0.6), Inches(6.3), Inches(8.0), Inches(0.8),
         "部署方式：Docker Compose 一键启动 5 个服务",
         ["mysql → redis → fastapi → frontend → metabase"],
         icon_text="🚀", icon_color=ACCENT_GREEN,
         bg_color=VERY_LIGHT_BLUE, border_color=ACCENT_GREEN)


# =============================================================================
# SLIDE 9: Dashboard (with screenshot)
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "07  核心功能 — 仪表板", font_size=28, bold=True, color=WHITE)

# Dashboard screenshot
dashboard_path = os.path.join(SCREENSHOTS, 'dashboard.png')
if os.path.exists(dashboard_path):
    slide.shapes.add_picture(dashboard_path, Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.0))

# KPI cards on right
kpi_items = [
    ("💰", "总销售额", "聚合所有订单金额", ACCENT_GREEN),
    ("📦", "总订单数", "统计订单总量", PRIMARY_BLUE),
    ("👥", "客户总数", "去重客户数量", ACCENT_ORANGE),
    ("📈", "环比增长率", "与上期数据对比", ACCENT_PURPLE),
]

add_textbox(slide, Inches(8.8), Inches(1.3), Inches(4), Inches(0.5),
            "KPI 指标卡片", font_size=18, bold=True, color=DARK_BLUE)

for i, (icon, title, desc, color) in enumerate(kpi_items):
    y = Inches(2.0) + Inches(i * 1.15)
    # KPI card
    card = add_shape(slide, Inches(8.8), y, Inches(4.0), Inches(1.0),
                     fill_color=WHITE, line_color=color, line_width=Pt(2))
    add_icon_circle(slide, Inches(9.0), y + Inches(0.2), Inches(0.55), title[0],
                    fill_color=color, font_size=18)
    add_textbox(slide, Inches(9.7), y + Inches(0.15), Inches(2.8), Inches(0.35),
                title, font_size=16, bold=True, color=color)
    add_textbox(slide, Inches(9.7), y + Inches(0.55), Inches(2.8), Inches(0.3),
                desc, font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 10: Charts
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "07  核心功能 — 可视化图表", font_size=28, bold=True, color=WHITE)

# Reports screenshot
reports_path = os.path.join(SCREENSHOTS, 'reports.png')
if os.path.exists(reports_path):
    slide.shapes.add_picture(reports_path, Inches(0.4), Inches(1.3), Inches(8.0), Inches(5.0))

# Chart types on right
chart_types = [
    ("📊", "销售趋势图", "折线图展示时间序列", LIGHT_BLUE),
    ("📈", "产品排行图", "柱状图展示 TOP N", ACCENT_GREEN),
    ("🥧", "品类占比图", "饼图展示构成比例", ACCENT_ORANGE),
    ("🌍", "地域分布图", "地图展示区域分布", ACCENT_PURPLE),
]

add_textbox(slide, Inches(8.8), Inches(1.3), Inches(4), Inches(0.5),
            "图表类型", font_size=18, bold=True, color=DARK_BLUE)

for i, (icon, title, desc, color) in enumerate(chart_types):
    y = Inches(2.0) + Inches(i * 1.15)
    card = add_shape(slide, Inches(8.8), y, Inches(4.0), Inches(1.0),
                     fill_color=WHITE, line_color=color, line_width=Pt(2))
    add_textbox(slide, Inches(9.1), y + Inches(0.15), Inches(3.5), Inches(0.35),
                title, font_size=16, bold=True, color=color)
    add_textbox(slide, Inches(9.1), y + Inches(0.55), Inches(3.5), Inches(0.3),
                desc, font_size=13, color=MEDIUM_GRAY)

# Bottom note
add_textbox(slide, Inches(8.8), Inches(6.8), Inches(4), Inches(0.4),
            "基于 ECharts 5.x 实现\n支持响应式布局", font_size=12, color=MEDIUM_GRAY)


# =============================================================================
# SLIDE 11: AI Smart Query
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "08  AI 智能问数 ⭐", font_size=28, bold=True, color=WHITE)

# AI query screenshot
ai_path = os.path.join(SCREENSHOTS, 'ai-query.png')
if os.path.exists(ai_path):
    slide.shapes.add_picture(ai_path, Inches(0.4), Inches(1.3), Inches(6.5), Inches(4.5))

# Process steps on right
add_textbox(slide, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.5),
            "系统处理流程", font_size=20, bold=True, color=DARK_BLUE)

process_steps = [
    ("1", "🧠", "AI 理解语义", "解析自然语言问题"),
    ("2", "💻", "自动生成 SQL", "基于百炼 Qwen 模型"),
    ("3", "▶️", "执行查询", "MySQL 数据库查询"),
    ("4", "📊", "展示结果", "表格/图表可视化"),
]
for i, (num, icon, title, desc) in enumerate(process_steps):
    y = Inches(2.0) + Inches(i * 1.0)
    add_icon_circle(slide, Inches(7.5), y, Inches(0.5), num, fill_color=PRIMARY_BLUE, font_size=16)
    add_textbox(slide, Inches(8.2), y - Inches(0.02), Inches(4), Inches(0.35),
                title, font_size=16, bold=True, color=DARK_BLUE)
    add_textbox(slide, Inches(8.2), y + Inches(0.35), Inches(4), Inches(0.3),
                desc, font_size=13, color=MEDIUM_GRAY)

# SQL example
add_shape(slide, Inches(7.3), Inches(6.1), Inches(5.5), Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(7.5), Inches(6.15), Inches(5), Inches(0.3),
            "生成 SQL 示例", font_size=14, bold=True, color=ACCENT_CYAN)
sql_text = 'SELECT p.product_name, SUM(oi.quantity * oi.price)\nFROM ods_order_items oi JOIN ods_products p ...\nWHERE o.order_date >= DATE_SUB(CURDATE(), INTERVAL 1 MONTH)\nORDER BY total_sales DESC LIMIT 1;'
add_textbox(slide, Inches(7.5), Inches(6.5), Inches(5), Inches(0.7),
            sql_text, font_size=11, color=RGBColor(0xCC, 0xD6, 0xDD))

# Tech metrics
add_textbox(slide, Inches(0.4), Inches(6.0), Inches(6.5), Inches(0.4),
            "技术指标：响应时间 < 3 秒  |  查询准确率 > 80%  |  支持多表 JOIN、聚合、时间表达式",
            font_size=14, bold=True, color=PRIMARY_BLUE)


# =============================================================================
# SLIDE 12: AI Query Examples Table
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "08  AI 智能问数 — 典型查询示例", font_size=28, bold=True, color=WHITE)

# Query examples table
query_table = [
    ["查询类型", "自然语言问题", "生成 SQL 特点"],
    ["简单查询", "\"有多少个产品？\"", "单表 COUNT"],
    ["聚合查询", "\"各品类销售占比\"", "GROUP BY + SUM"],
    ["排行查询", "\"销售额最高的产品\"", "ORDER BY + LIMIT"],
    ["时间查询", "\"上个月销售数据\"", "DATE_SUB 时间过滤"],
    ["多表查询", "\"哪个客户下单最多\"", "JOIN + GROUP BY"],
    ["复合查询", "\"华东区上月销量 TOP3\"", "多条件组合筛选"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(8.0), Inches(3.5), query_table,
          col_widths=[Inches(1.5), Inches(4.0), Inches(2.5)])

# Metabase section
add_textbox(slide, Inches(0.6), Inches(5.3), Inches(5), Inches(0.5),
            "Metabase BI 集成", font_size=20, bold=True, color=DARK_BLUE)

metabase_path = os.path.join(SCREENSHOTS, 'metabase.png')
if os.path.exists(metabase_path):
    slide.shapes.add_picture(metabase_path, Inches(0.6), Inches(5.8), Inches(5.5), Inches(1.4))

# Right side capabilities
add_card(slide, Inches(9.0), Inches(1.5), Inches(3.8), Inches(5.5),
         "双 BI 方案优势",
         [
             "拖拽式报表创建",
             "丰富的图表类型",
             "数据源管理",
             "报表导出分享",
             "自研前端 + Metabase 互补",
         ],
         icon_text="💡", icon_color=ACCENT_ORANGE,
         bg_color=VERY_LIGHT_BLUE, border_color=ACCENT_ORANGE)


# =============================================================================
# SLIDE 13: Test Results
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "09  系统测试与部署", font_size=28, bold=True, color=WHITE)

# Performance table
perf_table = [
    ["指标", "目标值", "实测值", "达成"],
    ["页面加载时间", "< 3 秒", "< 2 秒", "✅ 达成"],
    ["API 响应时间", "< 500ms", "< 200ms", "✅ 达成"],
    ["AI 查询响应", "< 5 秒", "< 3 秒", "✅ 达成"],
    ["数据准确性", "100%", "100%", "✅ 达成"],
    ["系统可用性", "99%", "99.9%", "✅ 达成"],
]
add_table(slide, Inches(0.6), Inches(1.5), Inches(7.0), Inches(3.0), perf_table,
          col_widths=[Inches(2.0), Inches(1.5), Inches(1.5), Inches(2.0)])

# Test coverage
add_textbox(slide, Inches(0.6), Inches(4.8), Inches(5), Inches(0.5),
            "测试覆盖", font_size=20, bold=True, color=DARK_BLUE)

test_items = [
    ("单元测试", "85% 覆盖率", ACCENT_GREEN),
    ("集成测试", "核心流程 100% 覆盖", PRIMARY_BLUE),
    ("压力测试", "支持 100+ 并发用户", ACCENT_ORANGE),
]
for i, (title, desc, color) in enumerate(test_items):
    y = Inches(5.4) + Inches(i * 0.55)
    add_icon_circle(slide, Inches(0.7), y, Inches(0.4), "✓", fill_color=color, font_size=14)
    add_textbox(slide, Inches(1.3), y - Inches(0.02), Inches(3), Inches(0.3),
                title, font_size=15, bold=True, color=color)
    add_textbox(slide, Inches(1.3), y + Inches(0.3), Inches(3), Inches(0.25),
                desc, font_size=12, color=MEDIUM_GRAY)

# Deployment card
add_card(slide, Inches(8.0), Inches(1.5), Inches(4.8), Inches(5.5),
         "Docker Compose 部署",
         [
             "一键启动：docker-compose up -d",
             "5 个服务：mysql / redis / fastapi / frontend / metabase",
             "数据持久化：MySQL volume 挂载",
             "日志收集：输出到 logs/ 目录",
             "环境隔离：.env 管理配置",
         ],
         icon_text="🚀", icon_color=ACCENT_GREEN,
         bg_color=VERY_LIGHT_BLUE, border_color=ACCENT_GREEN)


# =============================================================================
# SLIDE 14: Innovation Summary
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "10  创新点总结", font_size=28, bold=True, color=WHITE)

innovations = [
    ("1", "🏗️", "数仓分层架构", "经典四层设计（ODS→DWD→DWS→ADS）\n数据解耦，易于维护扩展\n共 4 层 24 张表，逐层加工", PRIMARY_BLUE),
    ("2", "🤖", "AI 智能问数", "自然语言转 SQL\n降低业务人员使用门槛\n响应 < 3 秒，准确率 > 80%", ACCENT_GREEN),
    ("3", "💡", "双 BI 方案", "自研前端 + Metabase\n灵活选择，优势互补\n固定报表 + 自助分析", ACCENT_ORANGE),
    ("4", "🚀", "一键部署", "Docker Compose 编排\n5 个服务快速上线\n降低运维成本", ACCENT_PURPLE),
]

for i, (num, icon, title, desc, color) in enumerate(innovations):
    x = Inches(0.6) + Inches(i * 3.15)
    # Card background
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(5.5), fill_color=WHITE,
              line_color=color, line_width=Pt(3))
    # Number circle
    add_icon_circle(slide, x + Inches(1.0), Inches(1.8), Inches(0.6), num, fill_color=color, font_size=22)
    # Title
    add_textbox(slide, x + Inches(0.2), Inches(2.6), Inches(2.5), Inches(0.5),
                title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.2), Inches(3.2), Inches(2.5), Inches(3.5),
                desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 15: Summary & Future Work
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), fill_color=DARK_BLUE)
add_textbox(slide, Inches(0.6), Inches(0.25), Inches(5), Inches(0.6),
            "11  总结与展望", font_size=28, bold=True, color=WHITE)

# Completed work
add_card(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.0),
         "✅ 已完成工作",
         [
             "完整的数仓分层架构（4 层 24 表）",
             "ETL 数据流水线（自动化调度）",
             "BI 可视化报表（自研 + Metabase）",
             "AI 智能问数功能（准确率 > 80%）",
         ],
         icon_text="🎯", icon_color=ACCENT_GREEN,
         bg_color=WHITE, border_color=ACCENT_GREEN)

# Future work
add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.4),
         "🔮 后续优化方向",
         [
             "真实 ERP 系统深度集成",
             "移动端应用开发",
             "实时数据处理（Kafka / Flink）",
             "多租户支持 + 更多 AI 能力",
         ],
         icon_text="🚀", icon_color=ACCENT_ORANGE,
         bg_color=WHITE, border_color=ACCENT_ORANGE)

# Right side - visual summary
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.7), fill_color=DARK_BLUE)
add_textbox(slide, Inches(7.5), Inches(2.0), Inches(4.8), Inches(0.6),
            "项目成果一览", font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(9.2), Inches(2.7), Inches(1.5), Inches(0.04), fill_color=ACCENT_CYAN)

成果 = [
    ("数仓表数", "24 张"),
    ("ETL 脚本", "4 层流水线"),
    ("图表类型", "4+ 种"),
    ("AI 准确率", "> 80%"),
    ("服务数量", "5 个"),
    ("部署方式", "Docker Compose"),
]
for i, (label, value) in enumerate(成果):
    col = i % 2
    row = i // 2
    x = Inches(7.8) + Inches(col * 2.8)
    y = Inches(3.1) + Inches(row * 1.1)
    add_textbox(slide, x, y, Inches(2.5), Inches(0.35),
                label, font_size=14, color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.35), Inches(2.5), Inches(0.45),
                value, font_size=24, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 16: Thank You
# =============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_gradient(slide)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), fill_color=ACCENT_CYAN)

add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.0),
            "致谢", font_size=42, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(5.8), Inches(2.9), Inches(1.7), Inches(0.05), fill_color=ACCENT_CYAN)

thanks_items = [
    "感谢导师的悉心指导",
    "感谢评委老师的聆听",
    "敬请批评指正",
]
for i, text in enumerate(thanks_items):
    y = Inches(3.4) + Inches(i * 0.7)
    add_textbox(slide, Inches(1.5), y, Inches(10.3), Inches(0.6),
                text, font_size=22, color=RGBColor(0xCC, 0xD6, 0xDD), alignment=PP_ALIGN.CENTER)

# Q&A
add_shape(slide, Inches(4.5), Inches(5.6), Inches(4.3), Inches(1.2),
          fill_color=PRIMARY_BLUE)
add_textbox(slide, Inches(4.5), Inches(5.75), Inches(4.3), Inches(0.5),
            "Q & A", font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.5), Inches(6.25), Inches(4.3), Inches(0.4),
            "欢迎提问", font_size=16, color=RGBColor(0xCC, 0xD6, 0xDD), alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(7.44), prs.slide_width, Inches(0.06), fill_color=PRIMARY_BLUE)


# ===== Save =====
output_path = os.path.join(BASE_DIR, 'ERP-BI 答辩_new.pptx')
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
